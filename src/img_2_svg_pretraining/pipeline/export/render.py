"""Render a multipage PDF into frames, video, and slides.

Consolidated from `export/pdf_to_video.py` and `export/tikz/ppt_convertor.py`,
whose algorithms are preserved -- notably the page-0 dimension locking in the
video path, which exists because a frame can grow beyond the first page's
size (e.g. a magnifier effect crossing the canvas edge) and a video writer
needs every frame identically sized.

Note this module rasterizes with its own `pdftoppm` call rather than
`viewer/compile.py::compile_tikz`, which passes `-singlefile` and therefore
only ever produces page 1 -- fine for a compile check, useless for frames.
"""
from __future__ import annotations

import subprocess
from pathlib import Path

DEFAULT_DPI = 300
DEFAULT_FPS = 2


class RenderError(Exception):
    pass


def compile_pdf(tex_source: str, out_pdf: Path, timeout: int = 300) -> Path:
    """Compile LaTeX to a (possibly multipage) PDF."""
    out_pdf = Path(out_pdf)
    work_dir = out_pdf.parent / f"_work_{out_pdf.stem}"
    work_dir.mkdir(parents=True, exist_ok=True)
    tex_file = work_dir / "doc.tex"
    tex_file.write_text(tex_source, encoding="utf-8")

    try:
        proc = subprocess.run(
            ["latexmk", "-pdf", "-interaction=nonstopmode", "-halt-on-error",
             f"-output-directory={work_dir}", str(tex_file)],
            cwd=work_dir, capture_output=True, text=True, timeout=timeout,
        )
        produced = work_dir / "doc.pdf"
        if proc.returncode != 0 or not produced.exists():
            log = (proc.stdout or "") + "\n" + (proc.stderr or "")
            raise RenderError(f"latexmk failed:\n{log[-4000:]}")
        out_pdf.parent.mkdir(parents=True, exist_ok=True)
        produced.replace(out_pdf)
        return out_pdf
    except subprocess.TimeoutExpired as e:
        raise RenderError(f"compilation timed out after {timeout}s: {e}") from e
    finally:
        import shutil
        shutil.rmtree(work_dir, ignore_errors=True)


def pdf_to_frames(pdf_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> list[Path]:
    """Rasterize every page to a PNG. No `-singlefile`, so all pages render."""
    pdf_path, out_dir = Path(pdf_path), Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    proc = subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(out_dir / "frame")],
        capture_output=True, text=True,
    )
    frames = sorted(out_dir.glob("frame*.png"))
    if proc.returncode != 0 or not frames:
        raise RenderError(f"pdftoppm failed: {proc.stdout}\n{proc.stderr}")
    return frames


def _load_uniform_frames(frames: list[Path]):
    """Load frames, resizing any that differ from the first.

    A page can render larger than page 1 when an effect overflows the canvas;
    video writers require constant dimensions, so later pages are scaled to
    match rather than dropped.
    """
    import cv2
    import numpy as np

    images = []
    target = None
    for path in frames:
        img = cv2.imread(str(path))
        if img is None:
            continue
        if target is None:
            target = (img.shape[1], img.shape[0])
        elif (img.shape[1], img.shape[0]) != target:
            img = cv2.resize(img, target, interpolation=cv2.INTER_AREA)
        images.append(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
    if not images:
        raise RenderError("no readable frames")
    return images


def frames_to_mp4(frames: list[Path], out_path: Path, fps: int = DEFAULT_FPS) -> Path:
    import imageio

    images = _load_uniform_frames(frames)
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    # macro_block_size=2 keeps libx264 from silently padding odd dimensions.
    with imageio.get_writer(str(out_path), fps=fps, codec="libx264",
                            macro_block_size=2) as writer:
        for image in images:
            writer.append_data(image)
    return out_path


def frames_to_gif(frames: list[Path], out_path: Path, fps: int = DEFAULT_FPS) -> Path:
    import imageio

    images = _load_uniform_frames(frames)
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    imageio.mimsave(str(out_path), images, duration=1.0 / max(fps, 1), loop=0)
    return out_path


def pdf_to_pptx(pdf_path: Path, out_path: Path, dpi: int = DEFAULT_DPI) -> Path:
    """One PDF page per slide, rendered in memory (no temp files)."""
    import io

    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    pdf_path, out_path = Path(pdf_path), Path(out_path)
    doc = fitz.open(str(pdf_path))
    try:
        if doc.page_count == 0:
            raise RenderError(f"{pdf_path} has no pages")

        presentation = Presentation()
        first = doc.load_page(0).rect
        emu_per_point = 914400 / 72
        presentation.slide_width = Emu(int(first.width * emu_per_point))
        presentation.slide_height = Emu(int(first.height * emu_per_point))
        blank = presentation.slide_layouts[6]

        zoom = dpi / 72
        for index in range(doc.page_count):
            pixmap = doc.load_page(index).get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            slide = presentation.slides.add_slide(blank)
            slide.shapes.add_picture(
                io.BytesIO(pixmap.tobytes("png")), 0, 0,
                width=presentation.slide_width, height=presentation.slide_height,
            )

        out_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(out_path))
        return out_path
    finally:
        doc.close()
