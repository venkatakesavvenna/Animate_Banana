import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
import os

def convert_pdf_to_pptx(pdf_path: str, ppt_path: str):
    """
    Converts a multi-page PDF into a PowerPoint presentation, 
    mapping each page to a single slide.
    """
    # 1. Initialize the PowerPoint presentation
    prs = Presentation()
    
    # 2. Select a completely blank slide layout (index 6)
    blank_slide_layout = prs.slide_layouts[6]
    
    # 3. Open the multi-page PDF
    print(f"Opening {pdf_path}...")
    doc = fitz.open(pdf_path)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # 4. Render the PDF page to a high-res image (300 DPI)
        pix = page.get_pixmap(dpi=300)
        
        # 5. Save the image to an in-memory buffer (avoids writing temp files)
        image_stream = io.BytesIO(pix.tobytes("png"))
        
        # 6. Dynamically resize the PPT canvas based on the first PDF page
        if page_num == 0:
            # PyMuPDF uses points (72 per inch). python-pptx uses EMUs (914400 per inch).
            ppt_width = int(page.rect.width * (914400 / 72))
            ppt_height = int(page.rect.height * (914400 / 72))
            prs.slide_width = ppt_width
            prs.slide_height = ppt_height
        
        # 7. Add a new slide and drop the image onto it
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            image_stream, 
            0, 0, 
            width=prs.slide_width, 
            height=prs.slide_height
        )
        
        print(f"Processed frame {page_num + 1}/{len(doc)}")
        
    # 8. Save the final presentation
    prs.save(ppt_path)
    print(f"Success! Saved presentation to: {ppt_path}")

if __name__ == "__main__":
    # ---------------------------------------------------------
    # HARDCODED CONFIGURATION
    # ---------------------------------------------------------
    input_pdf = "slide_bbox_new_export.pdf"
    
    # Generate the output filename dynamically
    name, _ = os.path.splitext(input_pdf)
    output_ppt = f"{name}.pptx"
    
    # ---------------------------------------------------------
    # EXECUTION
    # ---------------------------------------------------------
    try:
        convert_pdf_to_pptx(input_pdf, output_ppt)
    except FileNotFoundError:
        print(f"Error: Could not find '{input_pdf}'. Make sure the file exists.")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")