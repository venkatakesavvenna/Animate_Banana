import os
import math
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
import cv2

# --- Configuration ---
HTML_FILE_PATH = "alpha_making.html"
SLIDE_FRAMES_DIR = "temp_slides"
VIDEO_FRAMES_DIR = "temp_video"
OUTPUT_PPTX = "alpha_masking.pptx"
OUTPUT_MP4 = "alpha_masking.mp4"

# Tunables for the auto-fps heuristic (video only; slides don't use fps at all)
SAMPLES_PER_STAGE = 3   # how many video frames to sample within the shortest gap between stage transitions
FLOOR_FPS = 4           # never go below this even if the animation is very slow-paced
CEIL_FPS = 30           # never go above this even if the animation has very tight stage gaps

# Minimum gap (ms) between two stage-boundary times before we treat them as duplicates.
# Cheap pre-filter to avoid re-checking near-simultaneous keyframe times; the
# authoritative dedupe below compares each candidate's computed-style fingerprint,
# since two different timestamps can still resolve to the exact same state (e.g. an
# animation's own 0%-offset start, or a plateau's two shared-value edges).
SLIDE_DEDUPE_TOLERANCE_MS = 10

# Decimal places to round numeric computed-style values to before fingerprinting.
# Guards against floating-point jitter (e.g. 0.9999999 vs 1) being treated as a
# genuine change; does not affect the video pass at all.
FINGERPRINT_ROUND_DECIMALS = 2


def get_timeline_info(page):
    """
    Inspects every running CSS/Web Animation on the page and derives:
      - totalDurationMs: point at which the timeline is fully settled (max over finite
        animations' endTime, and over infinite animations' single-loop delay+duration,
        since 'duration' for an infinite loop means one full cycle)
      - minGapMs: smallest gap between any two stage-transition timestamps, used only
        to size FPS for the video so no stage gets skipped between video frames
    """
    return page.evaluate("""() => {
        const animations = document.getAnimations();
        if (animations.length === 0) {
            return { totalDurationMs: 0, minGapMs: 0 };
        }
        let maxEnd = 0;
        let eventTimes = [];
        animations.forEach(a => {
            const t = a.effect.getComputedTiming();
            const isInfinite = !isFinite(t.iterations);
            const end = isInfinite ? (t.delay + t.duration) : t.endTime;
            maxEnd = Math.max(maxEnd, end);
            eventTimes.push(t.delay, end);
        });
        const events = [...new Set(eventTimes)].sort((a, b) => a - b);
        let minGap = Infinity;
        for (let i = 1; i < events.length; i++) {
            minGap = Math.min(minGap, events[i] - events[i - 1]);
        }
        return {
            totalDurationMs: maxEnd,
            minGapMs: isFinite(minGap) ? minGap : maxEnd
        };
    }""")


def get_stage_boundary_times_ms(page):
    """
    Returns the sorted, deduped list of times (ms) at which some animation reaches
    a keyframe boundary -- i.e. a genuine visual "stage" (an element just finished
    fading/popping in, the sliding box just landed on a new target, etc).

    Uses each animation's own effect.getKeyframes() offsets (0..1) rather than
    parsing CSS text, so it works identically for delay-staged animations
    (fadeIn/popIn/colorPop) and for %-keyframe loops (the sliding box) without
    needing to know which family it is.
    """
    times = page.evaluate("""() => {
        const animations = document.getAnimations();
        let allTimes = [];
        animations.forEach(a => {
            const t = a.effect.getComputedTiming();
            const keyframes = a.effect.getKeyframes();
            const duration = isFinite(t.duration) ? t.duration : 0;
            keyframes.forEach(kf => {
                const offset = kf.computedOffset !== undefined ? kf.computedOffset : (kf.offset || 0);
                allTimes.push(t.delay + offset * duration);
            });
        });
        return [...new Set(allTimes)].sort((a, b) => a - b);
    }""")

    # Dedupe times that are within tolerance of each other (near-simultaneous stage settles)
    deduped = []
    for t in times:
        if not deduped or (t - deduped[-1]) > SLIDE_DEDUPE_TOLERANCE_MS:
            deduped.append(t)
    return deduped


def suggest_fps(min_gap_ms: float, samples_per_stage: int = SAMPLES_PER_STAGE,
                 floor_fps: int = FLOOR_FPS, ceil_fps: int = CEIL_FPS) -> int:
    """Pick video FPS from the animation's own granularity instead of a fixed guess."""
    if not min_gap_ms or min_gap_ms <= 0:
        return floor_fps
    fps = math.ceil((samples_per_stage * 1000) / min_gap_ms)
    return max(floor_fps, min(ceil_fps, fps))


def get_state_fingerprint(page, time_ms):
    """
    Seeks every animation to time_ms, then reads back the *actual resolved
    values* of exactly the CSS properties each animation's own keyframes
    declare (opacity, transform, filter, or animated SVG geometry properties
    like x/y/width/height for something like the sliding box) via
    getComputedStyle. Two timestamps are the same visual state iff this
    fingerprint string is identical -- no rendering, no pixel comparison,
    no threshold to tune.
    """
    return page.evaluate("""(t) => {
        const animations = document.getAnimations();
        animations.forEach(anim => {
            anim.pause();
            anim.currentTime = t;
        });

        const parts = [];
        animations.forEach((a, idx) => {
            const target = a.effect.target;
            const keyframes = a.effect.getKeyframes();
            const propNames = new Set();
            keyframes.forEach(kf => {
                Object.keys(kf).forEach(k => {
                    if (!['offset', 'easing', 'composite', 'computedOffset'].includes(k)) {
                        propNames.add(k);
                    }
                });
            });

            const style = getComputedStyle(target);
            [...propNames].sort().forEach(prop => {
                const cssProp = prop.replace(/[A-Z]/g, m => '-' + m.toLowerCase());
                let value = style.getPropertyValue(cssProp);
                if (!value) {
                    // Falls back to the raw attribute for SVG geometry
                    // properties (x/y/width/height/etc) some browsers
                    // don't surface through getComputedStyle.
                    value = target.getAttribute(prop) || '';
                }
                // Round any numeric-looking value to avoid float jitter
                // between renders being treated as a real change.
                const asNum = parseFloat(value);
                if (!isNaN(asNum) && String(asNum) === value.replace(/px$/, '')) {
                    value = asNum.toFixed(""" + str(FINGERPRINT_ROUND_DECIMALS) + """);
                }
                parts.push(idx + ':' + prop + '=' + value);
            });
        });
        return parts.join('|');
    }""", time_ms)


def _seek_and_screenshot(page, time_ms, out_path, width, height):
    """Pause all animations and jump to an exact timestamp, then screenshot."""
    page.evaluate("""(t) => {
        document.getAnimations().forEach(anim => {
            anim.pause();
            anim.currentTime = t;
        });
    }""", time_ms)
    page.screenshot(path=out_path, clip={"x": 0, "y": 0, "width": width, "height": height})


def open_page(playwright, html_path):
    """Launches the browser, loads the HTML, and returns (browser, page, width, height)."""
    browser = playwright.chromium.launch(headless=True)
    page = browser.new_page()

    abs_path = os.path.abspath(html_path)
    page.goto(f"file://{abs_path}")

    dimensions = page.evaluate("""() => {
        const svg = document.querySelector('svg');
        return {
            width: svg.viewBox.baseVal.width || 1000,
            height: svg.viewBox.baseVal.height || 260
        };
    }""")
    width = int(dimensions['width'])
    height = int(dimensions['height'])
    page.set_viewport_size({"width": width, "height": height})

    return browser, page, width, height


def extract_slide_frames():
    """
    One frame per animation *stage* (keyframe boundary), not per fixed time-step.
    This is what should feed the PPTX -- a handful of frames, each showing a
    genuinely new visual state, instead of a video-density frame dump.
    """
    if not os.path.exists(SLIDE_FRAMES_DIR):
        os.makedirs(SLIDE_FRAMES_DIR)

    kept_frames = []
    last_fingerprint = None

    with sync_playwright() as p:
        print("Launching headless browser (slides pass)...")
        browser, page, width, height = open_page(p, HTML_FILE_PATH)

        stage_times = get_stage_boundary_times_ms(page)

        if not stage_times:
            print("No animation keyframe stages detected -- nothing to build slides from.")
            browser.close()
            return [], width, height

        print(f"Checking {len(stage_times)} candidate stage times for distinct states...")

        for i, t_ms in enumerate(stage_times):
            fingerprint = get_state_fingerprint(page, t_ms)

            # Only screenshot (and keep) a candidate whose actual resolved
            # CSS/attribute state differs from the last one we kept -- this
            # is decided from the code's own computed values, not pixels.
            if fingerprint != last_fingerprint:
                final_path = os.path.join(SLIDE_FRAMES_DIR, f"slide_{len(kept_frames):03d}.png")
                page.screenshot(path=final_path, clip={"x": 0, "y": 0, "width": width, "height": height})
                kept_frames.append(final_path)
                last_fingerprint = fingerprint

            print(f"\rChecked {i + 1}/{len(stage_times)} candidates -> "
                  f"{len(kept_frames)} distinct slides kept", end="")

        browser.close()
        print(f"\nSlide frame extraction complete: {len(kept_frames)} distinct slides.")
        return kept_frames, width, height


def extract_video_frames():
    """
    Uniform time-step sampling across the full animation duration, at an
    auto-detected FPS -- this is what should feed the MP4, since a video wants
    smooth continuous motion rather than discrete stage snapshots.
    """
    if not os.path.exists(VIDEO_FRAMES_DIR):
        os.makedirs(VIDEO_FRAMES_DIR)

    frame_files = []

    with sync_playwright() as p:
        print("Launching headless browser (video pass)...")
        browser, page, width, height = open_page(p, HTML_FILE_PATH)

        timeline = get_timeline_info(page)
        animation_duration_ms = timeline["totalDurationMs"]
        fps = suggest_fps(timeline["minGapMs"])

        if animation_duration_ms <= 0:
            print("No running animations detected -- nothing to build a video from.")
            browser.close()
            return [], width, height, fps

        print(f"Detected total duration: {animation_duration_ms:.0f}ms "
              f"(smallest stage gap: {timeline['minGapMs']:.0f}ms) -> using {fps} FPS")

        total_frames = int((animation_duration_ms / 1000) * fps)
        time_step_ms = 1000 / fps

        print(f"Extracting {total_frames + 1} video frames...")

        for i in range(total_frames + 1):
            current_time = i * time_step_ms
            frame_path = os.path.join(VIDEO_FRAMES_DIR, f"frame_{i:04d}.png")
            _seek_and_screenshot(page, current_time, frame_path, width, height)
            frame_files.append(frame_path)
            print(f"\rCaptured frame {i}/{total_frames}", end="")

        browser.close()
        print("\nVideo frame extraction complete.")
        return frame_files, width, height, fps


def create_pptx(frame_files, width, height):
    """Compiles the stage-boundary slide frames into a PowerPoint presentation."""
    if not frame_files:
        print("No frames to build a PPTX from -- skipping.")
        return

    print("Generating PowerPoint...")
    prs = Presentation()

    width_in = width / 100.0
    height_in = height / 100.0
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    blank_slide_layout = prs.slide_layouts[6]  # fully blank layout

    for frame in frame_files:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(frame, 0, 0, width=Inches(width_in), height=Inches(height_in))

    prs.save(OUTPUT_PPTX)
    print(f"Saved: {OUTPUT_PPTX}")


def create_mp4(frame_files, width, height, fps):
    """Stitches the uniformly-sampled video frames into an MP4 at the detected FPS."""
    if not frame_files:
        print("No frames to build an MP4 from -- skipping.")
        return

    print("Generating MP4 Video...")
    fourcc = cv2.VideoWriter_fourcc(*'mp4v')
    out = cv2.VideoWriter(OUTPUT_MP4, fourcc, fps, (width, height))

    for frame in frame_files:
        img = cv2.imread(frame)
        out.write(img)

    out.release()
    print(f"Saved: {OUTPUT_MP4}")


if __name__ == "__main__":
    slide_frames, slide_w, slide_h = extract_slide_frames()
    create_pptx(slide_frames, slide_w, slide_h)

    video_frames, video_w, video_h, detected_fps = extract_video_frames()
    create_mp4(video_frames, video_w, video_h, detected_fps)

    print("All tasks completed successfully!")